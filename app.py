import os
import io
import json
import uuid
from flask import Flask, render_template, request, send_file
from flask_socketio import SocketIO, emit
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Initialize Flask app
app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('FLASK_SECRET_KEY')
socketio = SocketIO(app, cors_allowed_origins="*")

# Initialize Gemini API
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-1.5-flash')

# Store active presentation generation jobs
active_jobs = {}

@app.route('/')
def index():
    return render_template('index.html')

@socketio.on('connect')
def handle_connect():
    print('Client connected')

@socketio.on('disconnect')
def handle_disconnect():
    print('Client disconnected')

@socketio.on('generate_presentation')
def handle_generate_presentation(data):
    topic = data.get('topic', '')
    num_slides = int(data.get('num_slides', 5))
    
    # Validate number of slides
    num_slides = max(1, min(20, num_slides))
    
    include_images = data.get('include_images', False)
    style = data.get('style', 'professional')
    
    # Create a unique job ID
    job_id = str(uuid.uuid4())
    active_jobs[job_id] = {
        'status': 'processing',
        'progress': 0,
        'slides_data': []
    }
    
    emit('job_created', {'job_id': job_id})
    
    # Generate presentation content using Gemini
    try:
        emit('status_update', {'job_id': job_id, 'status': 'Generating content outline...'})
        
        # Modified prompt to generate actual content instead of instructions
        outline_prompt = f"""
        Create a detailed {num_slides}-slide presentation about "{topic}" in a {style} style.
        
        For each slide, provide:
        1. A clear and engaging slide title
        2. 3-5 bullet points of ACTUAL CONTENT (not meta-instructions)
        
        IMPORTANT: Each bullet point should contain real informative content that would appear on the slide, NOT meta-instructions about what kind of content to include.
        
        For example, instead of "Discuss the benefits of AI", provide the actual benefits like "Increases efficiency by automating repetitive tasks".
        
        Format the response as a JSON structure with this exact format:
        {{
            "presentation_title": "Title Here",
            "slides": [
                {{
                    "slide_title": "Slide 1 Title",
                    "bullet_points": ["Actual content point 1", "Actual content point 2", "Actual content point 3"]
                }},
                ...
            ]
        }}
        
        The first slide should be an introduction and the last slide should be a conclusion or summary.
        """
        
        outline_response = model.generate_content(outline_prompt)
        # Extract JSON from the response
        content_str = outline_response.text
        if '```json' in content_str:
            content_str = content_str.split('```json')[1].split('```')[0].strip()
        elif '```' in content_str:
            content_str = content_str.split('```')[1].split('```')[0].strip()
            
        # Parse JSON content
        try:
            content = json.loads(content_str)
        except json.JSONDecodeError:
            emit('error', {'message': 'Error parsing AI response'})
            return
            
        active_jobs[job_id]['slides_data'] = content
        
        # Generate the PowerPoint file
        presentation_path = generate_ppt(job_id, content, include_images)
        
        # Emit completion event
        emit('presentation_ready', {
            'job_id': job_id, 
            'download_url': f'/download/{job_id}'
        })
        
    except Exception as e:
        print(f"Error generating presentation: {e}")
        emit('error', {'message': f'Error generating presentation: {str(e)}'})

@app.route('/download/<job_id>')
def download_presentation(job_id):
    if job_id not in active_jobs:
        return "Presentation not found", 404
        
    # Create the presentation file
    job_data = active_jobs[job_id]
    prs = create_presentation(job_data['slides_data'], False)  # We already generated images if needed
    
    # Save to memory buffer
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    # Clean up job data after download
    filename = f"{job_data['slides_data']['presentation_title'].replace(' ', '_')}.pptx"
    
    # Schedule job cleanup (in a production app, you'd use a background task)
    def cleanup_job():
        if job_id in active_jobs:
            del active_jobs[job_id]
    
    # In a real app, use a proper background task here
    # For simplicity, we'll just delete immediately
    # threading.Timer(300, cleanup_job).start()
    
    return send_file(
        ppt_bytes,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

def generate_ppt(job_id, content_data, include_images=False):
    """Generate the PowerPoint presentation and return the path"""
    job_info = active_jobs[job_id]
    emit('status_update', {'job_id': job_id, 'status': 'Creating presentation slides...'})
    
    # Create and save the presentation
    prs = create_presentation(content_data, include_images)
    
    # Return job data
    job_info['status'] = 'complete'
    job_info['progress'] = 100
    emit('status_update', {'job_id': job_id, 'status': 'Presentation ready!', 'progress': 100})
    
    return job_id

def create_presentation(content_data, include_images=False):
    """Create a PowerPoint presentation from the content data"""
    prs = Presentation()
    
    # Set presentation properties
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content slide
    
    # Create title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = content_data['presentation_title']
    subtitle.text = "Generated with AI Presentation Generator"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add content slides
    total_slides = len(content_data['slides'])
    for i, slide_data in enumerate(content_data['slides']):
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = slide_data['slide_title']
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        # Add bullet points
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()  # Clear existing content
        
        for point in slide_data['bullet_points']:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(24)
        
        # If images are enabled, we would add them here
        # In a real implementation, you would generate or fetch relevant images
        if include_images and i > 0 and i < total_slides - 1:  # Skip adding images to first and last slides
            # This would be where you'd add image logic
            pass
    
    return prs

if __name__ == '__main__':
    socketio.run(app, debug=True)